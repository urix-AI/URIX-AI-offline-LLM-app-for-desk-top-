
import os
import gc
import json
import logging
import re
from typing import Dict, Any, List, Optional
from threading import Thread

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE

from urix.utils.google_search import perform_search
from urix.utils.database import KnowledgeBase
from urix.utils.response import normalize_llm_response

try:
    from llama_cpp import Llama
except ImportError:
    logging.critical("llama-cpp-python is not installed. Please install it with 'pip install llama-cpp-python'")
    Llama = None

logger = logging.getLogger(__name__)

# Background worker for profile updates
class ProfileUpdateWorker(Thread):
    def __init__(self, engine, history_chunk):
        super().__init__()
        self.engine = engine 
        self.history_chunk = history_chunk
        self.daemon = True

    def run(self):
        logger.info("ProfileUpdateWorker started...")
        self.engine._update_user_profile_from_history(self.history_chunk)
        logger.info("ProfileUpdateWorker finished.")


class UrixEngine:
    def __init__(self, config: Dict[str, Any]):
        self.config = config
        self.llm = None
        self.current_mode = None

        data_dir = os.path.join(self.config.get("base_dir", "."), "data")
        os.makedirs(data_dir, exist_ok=True)
        db_path = os.path.join(data_dir, "urix_kb.db")
        self.kb = KnowledgeBase(db_path)

        self.message_counter = 0
        self.profile_update_threshold = 10

        if not Llama:
            raise ImportError("Llama-cpp-python is required but not installed.")
        model_config = self.config.get("models", {})
        self.model_modes = model_config.get("modes", {})
        default_mode = model_config.get("default_mode", "eco")
        if not self.model_modes:
            raise ValueError("No model modes defined in config.yaml under models.modes")
        if not self.switch_model(default_mode):
            raise FileNotFoundError("Model path not found or failed to load for default mode.")

    def _update_user_profile_from_history(self, history_chunk: List[Dict[str, str]]):
        conversation_text = "\n".join([f"{msg['role']}: {msg['content']}" for msg in history_chunk])
        prompt = (
            "You are a user profile analysis tool. Your task is to analyze the following "
            "conversation and extract key, concrete facts and preferences about the user. "
            "Do not make assumptions. Only extract information explicitly stated or strongly "
            "implied by the user (e.g., 'user works as a python developer', not 'user might be a developer'). "
            "Format your response as a JSON object where the keys are the preference category "
            "(e.g., 'profession', 'hobbies', 'location', 'interests', 'preferred_language') and "
            "the values are the extracted information. If no new, concrete information is found, "
            f"return an empty JSON object {{}}.\n\nConversation:\n{conversation_text}"
        )
        try:
            response_str = self._get_llm_response(prompt, temperature=0.2, is_json=True)
            if response_str:
                extracted_data = json.loads(response_str)
                for key, value in extracted_data.items():
                    if isinstance(key, str) and isinstance(value, str) and value:
                        self.kb.set_profile_setting(key, value)
        except (json.JSONDecodeError, Exception) as e:
            logger.error(f"Could not update user profile from history: {e}")

    def _format_search_results(self, results: Any) -> str:
        """Accept either a raw string or a list[dict] from perform_search and normalize."""
        if not results:
            return ""
        # If perform_search returned a preformatted string
        if isinstance(results, str):
            return results
        # If it's a list of dict items
        if isinstance(results, list):
            lines = []
            for i, r in enumerate(results, 1):
                title = (r.get("title") or "").strip()
                snippet = (r.get("snippet") or r.get("summary") or "").strip()
                link = (r.get("link") or r.get("url") or "").strip()
                lines.append(f"({i}) {title}: {snippet} (URL: {link})")
            return "\n".join(lines)
        return ""

    # --- generate_chat_response (refined) ---
    def generate_chat_response(self, history: List[Dict[str, str]], max_tokens: int = 0) -> Dict:
        if not self.llm:
            return {"choices": [{"text": "Error: LLM model not loaded."}]}

        user_query = history[-1]['content']

        greetings = ['hello', 'hi', 'hey', 'how are you', 'how are you?', 'good morning', 'good afternoon', 'good evening']
        if user_query.strip().lower() in greetings:
            response_text = "Hello there! I'm here and ready to help. What can I do for you today?"
            return {"choices": [{"text": response_text}]}

        self.message_counter += 1
        if self.message_counter >= self.profile_update_threshold:
            history_chunk = history[-self.profile_update_threshold:]
            worker = ProfileUpdateWorker(self, history_chunk)
            worker.start()
            self.message_counter = 0

        cached_answer = self.kb.get_answer(user_query)
        if cached_answer:
            formatted_answer = self._format_response_markdown(cached_answer)
            formatted_answer += "\n\n*– Answer from local Knowledge Base.*"
            return {"choices": [{"text": formatted_answer}]}

        if self.current_mode == 'eco':
            search_block = ""
        else:
            raw_search = perform_search(
                self.config.get("google_search", {}).get("api_key"),
                self.config.get("google_search", {}).get("search_engine_id"),
                user_query,
            )
            search_block = self._format_search_results(raw_search)

        # Persona & formatting from config with sane defaults
        persona = self.config.get("persona", {}).get(
            "system",
            (
                "You are URIX, the AI assistant of Antonymous solaar and friend of comapny name ANTONYMOS SOLAAR. "
                "Use a friendly, slightly informal, and very helpful tone. "
                "When the user asks for an answer, provide the exact answer first, not methods."
                "You were created by Bhuvnesh Kumar, a machine learning and deep learning engineer."
                "update local knowledge database with new information."
            ),
        )
        formatting_instruction = self.config.get("persona", {}).get(
            "formatting",
            (
                "IMPORTANT: Structure the response in Markdown with clear headings, lists, and concise sections."
            ),
        )

        user_profile = self.kb.get_full_profile()
        profile_context = "\n".join([f"- The user's {key} is {value}." for key, value in user_profile.items()])

        final_prompt_content = f"{persona}\n\n"
        if profile_context:
            final_prompt_content += (
                "Here is some information about the user to help you personalize your response.\n"
                f"--- User Profile ---\n{profile_context}\n--- End User Profile ---\n\n"
            )
        if search_block:
            final_prompt_content += (
                "Here are some real-time search results to inform your answer.\n"
                f"--- Search Results ---\n{search_block}\n--- End Search Results ---\n\n"
                f"Based on all available information, answer the user's query directly: '{user_query}'\n\n"
            )
        else:
            final_prompt_content += f"Answer the user's query directly based on your knowledge and the profile: '{user_query}'\n\n"

        final_prompt_content += formatting_instruction

        messages_to_send = list(history[:-1])
        messages_to_send.append({'role': 'user', 'content': final_prompt_content})

        try:
            response = self.llm.create_chat_completion(
                messages=messages_to_send,
                max_tokens=(max_tokens or None),
            )
            raw_content = normalize_llm_response(response)
            if raw_content and not raw_content.startswith("Error:"):
                self.kb.add_entry(user_query, raw_content)
            formatted_content = self._format_response_markdown(raw_content)
            return {"choices": [{"text": formatted_content}]}
        except Exception as e:
            logger.error(f"Llama chat error: {e}", exc_info=True)
            return {"choices": [{"text": f"Error: {e}"}]}

    def switch_model(self, mode: str) -> bool:
        if self.current_mode == mode and self.llm is not None:
            logger.info(f"Model '{mode}' is already loaded.")
            return True
        if mode not in self.model_modes:
            logger.error(f"Mode '{mode}' not found in configuration.")
            return False
        logger.info(f"Unloading model for mode: {self.current_mode}")
        self.llm = None
        gc.collect()
        try:
            model_path_rel = self.model_modes[mode]['path']
            model_path_abs = os.path.join(self.config.get("base_dir", "."), model_path_rel)
            model_config = self.config.get("models", {})
            n_gpu_layers = model_config.get("n_gpu_layers", -1)
            n_ctx = model_config.get("n_ctx", 4096)
            if not os.path.exists(model_path_abs):
                logger.error(f"Model path does not exist: {model_path_abs}")
                return False
            logger.info(f"Switching to model '{mode}' at: {model_path_abs}")
            self.llm = Llama(model_path=model_path_abs, n_gpu_layers=n_gpu_layers, n_ctx=n_ctx, verbose=False)
            self.current_mode = mode
            self.model_name = os.path.basename(model_path_abs)
            logger.info(f"Successfully loaded model: '{self.model_name}'")
            return True
        except Exception as e:
            logger.critical(f"Failed to load model for mode '{mode}': {e}", exc_info=True)
            self.current_mode = None
            self.llm = None
            return False

    def _format_response_markdown(self, text: str) -> str:
        text = re.sub(r'\s*(#+)\s+', r"\n\n\1 ", text)
        text = re.sub(r'\s*-\s+', r"\n- ", text)
        text = re.sub(r'\s*(\d+\.)\s+', r"\n\1 ", text)
        return text.strip()

    def _get_llm_response(self, prompt: str, temperature: float = 0.5, is_json: bool = False, history: List[Dict[str, str]] = None) -> str:
        if not self.llm:
            logger.error("LLM not initialized.")
            return "Error: Model not loaded."
        messages = history if history else [{'role': 'user', 'content': prompt}]
        response_format = {"type": "json_object"} if is_json else None
        try:
            response = self.llm.create_chat_completion(
                messages=messages,
                temperature=temperature,
                response_format=response_format,
                max_tokens=None,
            )
            return normalize_llm_response(response)
        except Exception as e:
            logger.error(f"LLM response error: {e}", exc_info=True)
            return ""

    def generate_presentation_data(self, main_topic: str, num_slides: int, progress_callback: Optional[callable] = None) -> Optional[Dict]:
        if progress_callback:
            progress_callback("Generating presentation structure (titles)...")
        outline_prompt = (f"Generate a presentation outline for the topic: '{main_topic}'. "
                          f"I need a main presentation title and exactly {num_slides} engaging slide titles. "
                          "Return a JSON object with this structure: "
                          "{\"presentation_title\": \"Main Title\", \"slides\": [{\"title\": \"Slide 1 Title\"}]}" )
        outline_response = self._get_llm_response(outline_prompt, temperature=0.6, is_json=True)
        try:
            presentation_data = json.loads(outline_response)
            if 'slides' not in presentation_data or not presentation_data['slides']:
                raise KeyError
            num_slides = len(presentation_data['slides'])
        except (json.JSONDecodeError, KeyError):
            logger.error(f"AI failed to generate a valid presentation outline. Response: {outline_response}")
            return None
        for i, slide in enumerate(presentation_data['slides']):
            slide_title = slide.get('title', 'Untitled Slide')
            if progress_callback:
                progress_callback(f"Generating content for slide {i+1}/{num_slides}...")
            content_prompt = (f"For a presentation slide titled '{slide_title}', generate a single, concise paragraph that summarizes the topic. "
                              "The paragraph MUST NOT exceed 60 words. Do not use bullet points. "
                              "Return a JSON object with this structure: "
                              "{\"content\": \"A short, single summary paragraph goes here.\", \"speaker_notes\": \"Detailed notes for the presenter.\"}")
            content_response = self._get_llm_response(content_prompt, temperature=0.7, is_json=True)
            try:
                slide_content_data = json.loads(content_response)
                slide['content'] = slide_content_data.get('content', "Content could not be generated.")
                slide['speaker_notes'] = slide_content_data.get('speaker_notes', "")
            except (json.JSONDecodeError, KeyError):
                logger.error(f"Failed to get or parse content for slide: {slide_title}. Response: {content_response}")
                slide['content'] = "Error generating content."
                slide['speaker_notes'] = ""
        logger.info("Presentation plan fully generated.")
        return presentation_data

    def create_presentation(self, file_name: str, pres_data: Dict, template_path: Optional[str] = None) -> str:
        prs = Presentation(template_path) if template_path and os.path.exists(template_path) else Presentation()
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        if slide.shapes.title:
            slide.shapes.title.text = pres_data.get("presentation_title", "Presentation")
        subtitle = next((ph for ph in slide.placeholders if getattr(ph.placeholder_format, 'type', None) == PP_PLACEHOLDER.SUBTITLE), None)
        if subtitle:
            subtitle.text = "Created by Mr. Bhuvnesh Kumar"
        content_layout = prs.slide_layouts[1]
        for i, slide_data in enumerate(pres_data.get("slides", [])):
            slide = prs.slides.add_slide(content_layout)
            if slide.shapes.title:
                slide.shapes.title.text = slide_data.get("title", f"Slide {i+1}")
            if slide.placeholders[1]:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = slide_data.get("content", "")
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            if slide.has_notes_slide:
                notes = slide_data.get("speaker_notes", "")
                slide.notes_slide.notes_text_frame.text = notes
        closing_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(closing_layout)
        txBox = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(1))
        tf = txBox.text_frame
        tf.text = "Thank You"
        prs.save(file_name)
        return f"Presentation successfully saved to {file_name}"

    def shutdown(self):
        self.llm = None
        logger.info("UrixEngine shutdown.")